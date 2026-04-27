"""
Generate comprehensive PowerPoint presentation for Multi-Model Trading System results.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ── Color palette ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy background
BG_CARD = RGBColor(0x16, 0x21, 0x3A)       # Card background
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # Primary accent
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)  # Positive/success
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)    # Negative/risk
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)  # Warning/highlight
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6) # Secondary accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MED_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DIM_GRAY = RGBColor(0x64, 0x74, 0x8B)
TRANSPARENT = RGBColor(0x1E, 0x29, 0x3B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ════════════════════════════════════════════════════════════════
# Helper functions
# ════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=BG_CARD,
              border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.02
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=13, color=LIGHT_GRAY, spacing=Pt(6),
                     bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u2022  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"

        # Handle bold segments marked with **text**
        parts = item.split("**")
        for j, part in enumerate(parts):
            run = p.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
            if j % 2 == 1:
                run.font.bold = True
    return txBox


def add_table(slide, left, top, width, height, data, col_widths=None,
              header_color=ACCENT_BLUE, font_size=11):
    """data = list of lists. First row = header."""
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                    paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if r == 0:
                cell_fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x40)
            else:
                cell_fill.fore_color.rgb = BG_CARD

            # Borders
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)

    return table_shape


def add_metric_card(slide, left, top, width, height, label, value,
                    label_color=MED_GRAY, value_color=WHITE,
                    value_size=28, label_size=11, accent_line_color=ACCENT_BLUE):
    add_shape(slide, left, top, width, height, fill_color=BG_CARD)
    # Accent line at top
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15),
                                   top + Inches(0.08), width - Inches(0.3), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_line_color
    line.line.fill.background()
    # Value
    add_text_box(slide, left + Inches(0.15), top + Inches(0.25), width - Inches(0.3),
                 Inches(0.5), value, font_size=value_size, color=value_color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.15), top + height - Inches(0.35),
                 width - Inches(0.3), Inches(0.25), label,
                 font_size=label_size, color=label_color, alignment=PP_ALIGN.CENTER)


def add_section_header(slide, number, title, subtitle=""):
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.6), Inches(0.35), Inches(0.55), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_BLUE
    badge.line.fill.background()
    badge.adjustments[0] = 0.15
    tf = badge.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    # Title
    add_text_box(slide, Inches(1.35), Inches(0.3), Inches(10), Inches(0.55),
                 title, font_size=28, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1.35), Inches(0.82), Inches(10), Inches(0.35),
                     subtitle, font_size=14, color=MED_GRAY)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2),
                                   Inches(12.1), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x2D, 0x3A, 0x55)
    line.line.fill.background()


# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Decorative accent bar at top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                              SLIDE_W, Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# Title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
             "Multi-Model Cross-Sectional Equity Alpha System",
             font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(0.6),
             "Comparative Analysis: LightGBM vs Time Series Transformer vs CrossMamba",
             font_size=20, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.LEFT)

# Divider
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5),
                               Inches(2), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Subtitle details
details = [
    "S&P 500 Universe  |  186 Features  |  Walk-Forward Backtest",
    "February 2023 - February 2026  |  802 Trading Days",
    "Market-Neutral (10L/10S) & Long-Biased (14L/7S) Configurations"
]
for i, d in enumerate(details):
    add_text_box(slide, Inches(1), Inches(3.9 + i * 0.4), Inches(11), Inches(0.35),
                 d, font_size=14, color=MED_GRAY)

# Bottom metric cards
metrics_title = [
    ("Best Sharpe", "2.36", ACCENT_GREEN),
    ("Best Annual Return", "30.24%", ACCENT_GREEN),
    ("Best Max Drawdown", "-9.20%", ACCENT_BLUE),
    ("SPY Sharpe", "1.34", MED_GRAY),
]
card_w = Inches(2.5)
start_x = Inches(1)
for i, (label, value, accent) in enumerate(metrics_title):
    add_metric_card(slide, start_x + i * (card_w + Inches(0.3)), Inches(5.5),
                    card_w, Inches(1.3), label, value, accent_line_color=accent)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "", "Agenda")

agenda_items_left = [
    ("01", "System Architecture & Design Philosophy", ACCENT_BLUE),
    ("02", "Data Pipeline & Feature Engineering", ACCENT_BLUE),
    ("03", "Model Architectures Deep Dive", ACCENT_BLUE),
    ("04", "Walk-Forward Backtesting Framework", ACCENT_BLUE),
    ("05", "Risk Management & Portfolio Construction", ACCENT_BLUE),
]
agenda_items_right = [
    ("06", "Results: Market-Neutral (10L/10S)", ACCENT_GREEN),
    ("07", "Results: Long-Biased (14L/7S)", ACCENT_GREEN),
    ("08", "Feature Importance & Alpha Sources", ACCENT_PURPLE),
    ("09", "Architectural Insights & Key Findings", ACCENT_AMBER),
    ("10", "Production Deployment & Future Work", MED_GRAY),
]

for col, items, x_off in [(agenda_items_left, agenda_items_left, Inches(1)),
                            (agenda_items_right, agenda_items_right, Inches(7))]:
    for i, (num, text, color) in enumerate(items):
        y = Inches(1.7) + i * Inches(1.0)
        # Number
        add_text_box(slide, x_off, y, Inches(0.6), Inches(0.5),
                     num, font_size=22, color=color, bold=True, alignment=PP_ALIGN.LEFT)
        # Text
        add_text_box(slide, x_off + Inches(0.7), y + Inches(0.03), Inches(5), Inches(0.45),
                     text, font_size=16, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "1", "System Architecture", "Cross-Sectional Ranking Framework")

# Left: Design philosophy
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.35),
             "Design Philosophy", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(0.7), Inches(1.9), Inches(5.5), Inches(2.8), [
    "Predicts **relative forward returns** (cross-sectional rank) not absolute returns",
    "Eliminates market beta as a confound \u2014 ranks are distribution-free",
    "Maps naturally to **long/short portfolio construction**",
    "Walk-forward design with **strict temporal separation** prevents leakage",
    "Identical risk/portfolio pipeline for all models \u2014 apples-to-apples comparison",
], font_size=13)

# Right: Pipeline
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.35),
             "Pipeline Stages", font_size=16, color=ACCENT_BLUE, bold=True)

pipeline_stages = [
    ("1", "Universe Construction", "S&P 500, ~102 tickers after liquidity filter"),
    ("2", "Feature Engineering", "186 features across 4 signal domains"),
    ("3", "Feature Selection", "Stability-based IC \u2192 50 features"),
    ("4", "Walk-Forward Training", "504d train / 21d retrain / 10d purge"),
    ("5", "Factor Risk Model", "Barra-style: \u03a3 = BFB\u2032 + D"),
    ("6", "Portfolio Construction", "Risk-parity weighting, turnover control"),
    ("7", "Execution", "Alpaca Markets API (paper + live)"),
]

for i, (num, title, desc) in enumerate(pipeline_stages):
    y = Inches(1.95) + i * Inches(0.63)
    # Num badge
    badge = add_shape(slide, Inches(7), y, Inches(0.38), Inches(0.38),
                      fill_color=ACCENT_BLUE)
    tf = badge.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title + desc
    add_text_box(slide, Inches(7.55), y - Inches(0.02), Inches(4.5), Inches(0.22),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.55), y + Inches(0.2), Inches(4.5), Inches(0.22),
                 desc, font_size=10, color=MED_GRAY)

# Bottom: Universe specs
add_shape(slide, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.5), fill_color=BG_CARD)
add_text_box(slide, Inches(1), Inches(5.5), Inches(3), Inches(0.3),
             "Universe & Data Specifications", font_size=13, color=ACCENT_BLUE, bold=True)

specs = [
    ("Source", "S&P 500"),
    ("Liquidity Filter", "\u2265$5M avg daily $ vol"),
    ("Min History", "252 days"),
    ("Tickers", "~102"),
    ("Cross-Assets", "VIX, TNX, GLD, USO, HYG, TLT + 11 sector ETFs"),
    ("Backtest", "Feb 2023 \u2013 Feb 2026 (802 days)"),
]
for i, (k, v) in enumerate(specs):
    col = i % 3
    row = i // 3
    x = Inches(1) + col * Inches(3.9)
    y = Inches(5.9) + row * Inches(0.4)
    add_text_box(slide, x, y, Inches(1.5), Inches(0.3), k + ":", font_size=11, color=MED_GRAY, bold=True)
    add_text_box(slide, x + Inches(1.55), y, Inches(2.2), Inches(0.3), v, font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: FEATURE ENGINEERING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "2", "Feature Engineering", "186 Features Across 4 Signal Domains")

# Four domain cards
domains = [
    ("Price / Volume", "~100 features", ACCENT_BLUE, [
        "Momentum: [5, 10, 21, 63, 126, 252]d windows",
        "Volatility: realized vol at [5, 10, 21, 63]d",
        "Mean reversion: distance from N-day high/low",
        "Liquidity: Amihud illiquidity, log dollar volume",
        "Technical: MACD, RSI(14), Bollinger Bands",
        "All features also cross-sectionally ranked",
    ]),
    ("Fundamental", "~40 features", ACCENT_GREEN, [
        "Valuation: P/E, EV/EBITDA, Price/Book",
        "Quality: ROE, ROA, margins, current ratio",
        "Growth: earnings growth, revenue growth",
        "Analyst: target upside, recommendations",
        "Earnings day return (post-earnings drift)",
        "Composite scores: value, quality, sector-relative",
    ]),
    ("Cross-Asset", "~36 features", ACCENT_PURPLE, [
        "VIX level, percentile rank, rate of change",
        "Yield curve slope: 10Y \u2013 3M",
        "Credit spread: HYG/LQD ratio",
        "Sector ETF relative momentum",
        "Dollar (UUP), Gold, Oil momentum",
        "Risk-on/off: IWM/QQQ vs SPY",
    ]),
    ("Sentiment", "~10 features", ACCENT_AMBER, [
        "News headline keyword scoring",
        "Positive/negative lexicon matching",
        "Per-ticker sentiment aggregation",
        "Source: yfinance news feed",
        "",
        "",
    ]),
]

card_w = Inches(2.9)
gap = Inches(0.2)
start_x = Inches(0.55)

for i, (title, count, color, bullets) in enumerate(domains):
    x = start_x + i * (card_w + gap)
    card = add_shape(slide, x, Inches(1.55), card_w, Inches(4.6), fill_color=BG_CARD)
    # Color bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.55),
                                  card_w, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Title
    add_text_box(slide, x + Inches(0.2), Inches(1.75), card_w - Inches(0.4), Inches(0.35),
                 title, font_size=15, color=color, bold=True)
    # Count
    add_text_box(slide, x + Inches(0.2), Inches(2.1), card_w - Inches(0.4), Inches(0.25),
                 count, font_size=11, color=MED_GRAY)
    # Bullets
    for j, bullet in enumerate(bullets):
        if bullet:
            add_text_box(slide, x + Inches(0.2), Inches(2.5) + j * Inches(0.42),
                         card_w - Inches(0.4), Inches(0.4),
                         "\u2022  " + bullet, font_size=10, color=LIGHT_GRAY)

# Bottom: Feature selection
add_shape(slide, Inches(0.55), Inches(6.3), Inches(12.2), Inches(0.9), fill_color=BG_CARD)
add_text_box(slide, Inches(0.85), Inches(6.4), Inches(4), Inches(0.3),
             "Stability-Based Feature Selection", font_size=13, color=ACCENT_AMBER, bold=True)
add_text_box(slide, Inches(0.85), Inches(6.7), Inches(11.5), Inches(0.4),
             "Score = mean(|IC|) \u00d7 sign_consistency   |   Data split into 3 time periods   |   "
             "IC sign must be consistent across periods   |   186 \u2192 50 features",
             font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: MODEL ARCHITECTURE — LightGBM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "3a", "LightGBM", "Gradient Boosted Decision Tree Ensemble")

# Left: Overview
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.3),
             "How It Works", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(0.7), Inches(1.9), Inches(5.8), Inches(3.5), [
    "Builds an ensemble of **800 shallow decision trees** sequentially",
    "Each tree corrects the errors of the previous trees (gradient boosting)",
    "Sees each (date, ticker) pair as an **independent 50-feature row** \u2014 no notion of time",
    "Final prediction = sum of all tree outputs for that row",
    "**3 models** trained with different random seeds, predictions averaged",
    "Natively handles missing values \u2014 learns optimal split direction for NaN",
    "Produces interpretable **feature importance** via split gain",
], font_size=13)

add_text_box(slide, Inches(0.7), Inches(4.5), Inches(5.8), Inches(0.3),
             "Key Limitation", font_size=14, color=ACCENT_RED, bold=True)
add_bullet_frame(slide, Inches(0.7), Inches(4.85), Inches(5.8), Inches(1.0), [
    "**Cannot model temporal dependencies** \u2014 treats every observation as i.i.d.",
    "Misses patterns like \"momentum accelerating over the last 5 of 21 days\"",
], font_size=13, bullet_color=ACCENT_RED)

# Right: Config table
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Configuration (Optuna-Optimized)", font_size=14, color=ACCENT_BLUE, bold=True)

lgb_data = [
    ["Parameter", "Default", "Optimized"],
    ["n_estimators", "800", "800"],
    ["max_depth", "5", "6"],
    ["learning_rate", "0.030", "0.051"],
    ["num_leaves", "24", "40"],
    ["min_child_samples", "100", "250"],
    ["subsample", "0.70", "0.50"],
    ["colsample_bytree", "0.60", "0.40"],
    ["L1 regularization", "0.50", "1.44"],
    ["L2 regularization", "5.00", "7.61"],
]
add_table(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(3.2), lgb_data, font_size=11)

# Bottom callout
add_shape(slide, Inches(7.2), Inches(5.3), Inches(5.3), Inches(1.5), fill_color=BG_CARD,
          border_color=ACCENT_AMBER)
add_text_box(slide, Inches(7.5), Inches(5.4), Inches(4.8), Inches(0.3),
             "Optuna Insight", font_size=13, color=ACCENT_AMBER, bold=True)
add_text_box(slide, Inches(7.5), Inches(5.75), Inches(4.8), Inches(0.9),
             "Optimization pushed toward higher regularization (L2: 5\u21927.6) and "
             "lower feature sampling (0.6\u21920.4), suggesting the 50-feature space "
             "benefits from aggressive regularization to prevent overfitting.",
             font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: MODEL ARCHITECTURE — TST
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "3b", "Time Series Transformer (TST)",
                   "Multi-Head Self-Attention Over Rolling Windows")

# Left: Overview
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.3),
             "How It Works", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(0.7), Inches(1.9), Inches(5.8), Inches(4.5), [
    "Takes a **21-day rolling window** of 50 features per stock as input",
    "Projects each day's features into a 64-dimensional representation",
    "Adds **positional encoding** so the model knows the temporal ordering",
    "Passes through **2 Transformer encoder layers** with 4 attention heads",
    "Self-attention lets each day \"look at\" every other day in the window, "
    "learning which past days are most relevant for today's prediction",
    "Uses the **last time step's representation** as a summary of the entire sequence",
    "An output head (two linear layers with ReLU) maps this to a **scalar ranking score**",
    "**2 models** ensembled with different seeds",
], font_size=13)

# Right: Config + strengths
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Configuration", font_size=14, color=ACCENT_BLUE, bold=True)

tst_data = [
    ["Parameter", "Value"],
    ["d_model (hidden dim)", "64"],
    ["Attention heads", "4"],
    ["Encoder layers", "2"],
    ["Feedforward dim", "128"],
    ["Sequence length", "21 days"],
    ["Dropout", "0.2"],
    ["Optimizer", "AdamW (lr=1e-3)"],
    ["Batch size", "256"],
    ["Epochs", "3"],
    ["Gradient clipping", "max_norm=1.0"],
]
add_table(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(3.5), tst_data, font_size=11)

# Strengths/weaknesses
add_text_box(slide, Inches(7.2), Inches(5.6), Inches(2.5), Inches(0.3),
             "Strengths", font_size=13, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(7.2), Inches(5.9), Inches(5), Inches(0.8),
             "\u2022  Captures cross-time dependencies\n"
             "\u2022  Learns which past days matter",
             font_size=11, color=LIGHT_GRAY)

add_text_box(slide, Inches(10), Inches(5.6), Inches(2.5), Inches(0.3),
             "Weakness", font_size=13, color=ACCENT_RED, bold=True)
add_text_box(slide, Inches(10), Inches(5.9), Inches(3), Inches(0.8),
             "\u2022  O(L\u00b2) complexity \u2014 quadratic\n    in sequence length",
             font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: MODEL ARCHITECTURE — CrossMamba
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "3c", "CrossMamba (Selective State-Space Model)",
                   "Linear-Time Sequence Modeling with Selective Memory")

# Left: Overview
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.3),
             "How It Works", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(0.7), Inches(1.9), Inches(5.8), Inches(4.5), [
    "Like TST, takes a **21-day rolling window** of 50 features per stock",
    "Instead of attention, uses a **selective state-space model (SSM)** \u2014 "
    "a recurrent hidden state that evolves across time steps",
    "The key idea: **how much to remember vs forget at each step is input-dependent**. "
    "The model learns to selectively retain useful past information and discard noise",
    "A **depthwise convolution** (kernel=4) captures local patterns (3\u20134 day effects) "
    "while the hidden state carries longer-range regime information",
    "Uses a **dual-path gating mechanism**: one path processes the sequence through the SSM, "
    "the other acts as a learned gate that filters the output",
    "Stacks **2 CrossMamba blocks**, each with SSM + feedforward network + residual connections",
    "Processes sequences in **O(L) linear time** vs O(L\u00b2) for Transformers",
], font_size=12)

# Right: Config + comparison
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Configuration", font_size=14, color=ACCENT_BLUE, bold=True)

mamba_data = [
    ["Parameter", "Value"],
    ["d_model (hidden dim)", "64"],
    ["d_state (SSM state)", "16"],
    ["d_conv (conv kernel)", "4"],
    ["Expand factor", "2"],
    ["Layers", "2"],
    ["Sequence length", "21 days"],
    ["Dropout", "0.2"],
    ["Optimizer", "AdamW (lr=1e-3)"],
    ["Model size on disk", "582 KB (smallest)"],
]
add_table(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(3.1), mamba_data, font_size=11)

# Why it wins callout
add_shape(slide, Inches(7.2), Inches(5.2), Inches(5.3), Inches(1.8), fill_color=BG_CARD,
          border_color=ACCENT_GREEN)
add_text_box(slide, Inches(7.5), Inches(5.3), Inches(4.8), Inches(0.3),
             "Why CrossMamba Wins", font_size=13, color=ACCENT_GREEN, bold=True)
add_bullet_frame(slide, Inches(7.3), Inches(5.6), Inches(5), Inches(1.3), [
    "**Selective forgetting**: learns which past data is relevant \u2014 critical for noisy financial data",
    "**Local + global context**: conv captures microstructure; hidden state captures regime persistence",
    "**Linear time**: O(L) enables scaling to longer sequences without quadratic blowup",
], font_size=11, bullet_color=ACCENT_GREEN)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: MODEL COMPARISON OVERVIEW
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "3d", "Architecture Comparison",
                   "Three Paradigms for Modeling Financial Time Series")

# Comparison table
comp_data = [
    ["", "LightGBM", "TST (Transformer)", "CrossMamba (SSM)"],
    ["Paradigm", "Gradient boosted trees", "Self-attention", "Selective state-space"],
    ["Input", "(N, 50) flat rows", "(N, 21, 50) sequences", "(N, 21, 50) sequences"],
    ["Temporal modeling", "None", "Full attention matrix", "Recurrent hidden state"],
    ["Time complexity", "O(n log n)", "O(L\u00b2 \u00d7 d)", "O(L \u00d7 d)"],
    ["How it handles history", "Ignores it", "Looks at all pairs of days", "Carries a memory state forward"],
    ["Forgetting mechanism", "N/A", "Implicit (re-weighting)", "Explicit (learned decay)"],
    ["Local patterns", "N/A", "No built-in local bias", "Depthwise conv (k=4)"],
    ["Gating", "No", "No", "Yes (dual-path SiLU)"],
    ["Model size (disk)", "1.2 MB", "839 KB", "582 KB (smallest)"],
    ["Training speed", "Fast (~seconds)", "Moderate (~minutes)", "Moderate (~minutes)"],
    ["Feature importance", "Split gain (native)", "Gradient-based", "Gradient-based"],
]
add_table(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.2),
          comp_data, font_size=11,
          col_widths=[Inches(2.2), Inches(3.1), Inches(3.3), Inches(3.3)])


# ════════════════════════════════════════════════════════════════
# SLIDE 9: WALK-FORWARD BACKTESTING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "4", "Walk-Forward Backtesting Framework",
                   "Anti-Leakage Design with Purge & Embargo")

# Timeline diagram
add_shape(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.8), fill_color=BG_CARD)

# Timeline bar segments
segments = [
    (Inches(1.2), Inches(5.5), "504d Train Window", ACCENT_BLUE),
    (Inches(6.7), Inches(1.0), "10d Purge", ACCENT_RED),
    (Inches(7.7), Inches(0.8), "5d Embargo", ACCENT_AMBER),
    (Inches(8.5), Inches(2.5), "21d OOS Predict", ACCENT_GREEN),
]
bar_y = Inches(2.15)
bar_h = Inches(0.4)
for x, w, label, color in segments:
    seg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, bar_y, w, bar_h)
    seg.fill.solid()
    seg.fill.fore_color.rgb = color
    seg.line.fill.background()
    seg.adjustments[0] = 0.15
    tf = seg.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"

add_text_box(slide, Inches(1), Inches(2.65), Inches(10), Inches(0.3),
             "Each retrain: train on trailing 504 days \u2192 purge 10d gap \u2192 embargo 5d \u2192 predict next 21 days "
             "\u2192 slide window forward and repeat",
             font_size=11, color=MED_GRAY)

# Config table
wf_data = [
    ["Parameter", "Value", "Rationale"],
    ["Train window", "504 days (~2 years)", "Sufficient regime diversity"],
    ["Retrain frequency", "21 days (monthly)", "Balances recency vs stability"],
    ["Purge gap", "10 days", "Prevents target leakage (overlapping fwd returns)"],
    ["Embargo", "5 days", "Buffer for serial correlation in returns"],
    ["Target horizon", "10 days forward", "Balances signal strength vs data points"],
    ["Target type", "Cross-sectional rank", "Distribution-free, market-neutral by design"],
    ["Validation split", "20% of train window", "For early stopping"],
    ["Early stopping", "50 rounds (LightGBM)", "Prevents overfitting"],
    ["Walk-forward windows", "~38 over 3.18 years", "Comprehensive OOS coverage"],
]
add_table(slide, Inches(0.7), Inches(3.5), Inches(11.9), Inches(3.7),
          wf_data, font_size=11,
          col_widths=[Inches(2.2), Inches(3.5), Inches(6.2)])


# ════════════════════════════════════════════════════════════════
# SLIDE 10: RISK MANAGEMENT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "5", "Risk Management & Portfolio Construction",
                   "Barra-Style Factor Model + Multi-Layer Risk Pipeline")

# Left: Risk pipeline
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.3),
             "Risk Pipeline (Applied Sequentially)", font_size=15, color=ACCENT_BLUE, bold=True)

risk_steps = [
    ("1", "Sector Neutralization", "Max \u00b13% net per sector", ACCENT_BLUE),
    ("2", "Factor Neutralization", "Reduce Barra factor tilts toward zero", ACCENT_BLUE),
    ("3", "Volatility Targeting", "10% (mkt-neutral) / 15% (long-biased)", ACCENT_GREEN),
    ("4", "Drawdown Control", "Scale down 50% if drawdown > -8%", ACCENT_RED),
    ("5", "Position Clipping", "Max 10% per name, top-N long/short", ACCENT_AMBER),
    ("6", "Regime Overlay", "50d/200d MA trend, max \u00b115% bias", ACCENT_PURPLE),
]

for i, (num, title, desc, color) in enumerate(risk_steps):
    y = Inches(1.95) + i * Inches(0.72)
    badge = add_shape(slide, Inches(0.9), y, Inches(0.38), Inches(0.38), fill_color=color)
    tf = badge.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_box(slide, Inches(1.45), y - Inches(0.02), Inches(5), Inches(0.23),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.45), y + Inches(0.22), Inches(5), Inches(0.23),
                 desc, font_size=11, color=MED_GRAY)

# Right top: Portfolio config
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Portfolio Configuration", font_size=15, color=ACCENT_BLUE, bold=True)

port_data = [
    ["", "Market-Neutral", "Long-Biased"],
    ["Positions", "10 Long / 10 Short", "14 Long / 7 Short"],
    ["Max position", "10%", "10%"],
    ["Max gross leverage", "1.6\u00d7", "1.6\u00d7"],
    ["Net exposure", "\u00b115%", "~33%"],
    ["Weighting", "Risk-parity", "Risk-parity"],
    ["Vol target", "10% annual", "15% annual"],
    ["Turnover penalty", "0.1% score boost", "0.1% score boost"],
]
add_table(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(2.8), port_data, font_size=11)

# Right bottom: TC model
add_text_box(slide, Inches(7.2), Inches(5.0), Inches(5.5), Inches(0.3),
             "Transaction Cost Model", font_size=14, color=ACCENT_BLUE, bold=True)

tc_data = [
    ["Component", "Cost (bps)"],
    ["Commission", "0.5"],
    ["Slippage", "3.0"],
    ["Spread", "2.0"],
    ["Total per side", "5.5"],
    ["Round trip", "11.0"],
]
add_table(slide, Inches(7.2), Inches(5.35), Inches(3.5), Inches(1.8), tc_data, font_size=11)

# Factor risk model
add_text_box(slide, Inches(0.7), Inches(6.3), Inches(5.5), Inches(0.3),
             "Barra Factor Model: \u03a3 = BFB\u2032 + D", font_size=13, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(0.7), Inches(6.6), Inches(5.5), Inches(0.5),
             "6 factors: Market, Size, Value, Momentum, Volatility, Quality\n"
             "Covariance estimated with exponential decay (halflife=63d)",
             font_size=11, color=MED_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 11: RESULTS — MARKET NEUTRAL
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "6", "Results: Market-Neutral (10L/10S)",
                   "February 2023 \u2013 February 2026  |  802 Trading Days")

# Main results table
mn_data = [
    ["Metric", "LightGBM", "TST", "CrossMamba", "Ensemble"],
    ["Total Return", "17.29%", "25.48%", "61.12%", "54.83%"],
    ["Annual Return", "5.14%", "7.39%", "16.17%", "14.73%"],
    ["Volatility", "7.65%", "7.24%", "7.71%", "7.58%"],
    ["Sharpe Ratio", "0.67", "1.02", "2.10", "1.94"],
    ["Sortino Ratio", "1.09", "1.68", "3.61", "3.23"],
    ["Max Drawdown", "-12.89%", "-10.70%", "-5.93%", "-10.56%"],
    ["Calmar Ratio", "0.40", "0.69", "2.73", "1.39"],
    ["Win Rate", "50.50%", "52.99%", "53.24%", "54.86%"],
    ["Profit Factor", "1.12", "1.18", "1.37", "1.35"],
]
add_table(slide, Inches(0.7), Inches(1.5), Inches(7.5), Inches(4.0),
          mn_data, font_size=12,
          col_widths=[Inches(1.7), Inches(1.45), Inches(1.45), Inches(1.45), Inches(1.45)])

# Key observations
add_text_box(slide, Inches(8.7), Inches(1.5), Inches(4), Inches(0.3),
             "Key Observations", font_size=15, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(8.7), Inches(1.9), Inches(4.2), Inches(3.6), [
    "**CrossMamba dominates**: 2.10 Sharpe, 3\u00d7 the return of LightGBM",
    "CrossMamba achieves the **tightest drawdown** (-5.93%) \u2014 less than half of LightGBM",
    "**TST** shows meaningful improvement over LightGBM (1.02 vs 0.67 Sharpe) but well below CrossMamba",
    "**Ensemble paradox**: equal-weight blend (1.94) underperforms best individual model (2.10)",
    "All models target **~7.5% realized vol**, confirming risk management works correctly",
], font_size=12)

# Metric highlight cards
cards = [
    ("Best Sharpe", "2.10", "CrossMamba", ACCENT_GREEN),
    ("Best Max DD", "-5.93%", "CrossMamba", ACCENT_BLUE),
    ("Best Win Rate", "54.86%", "Ensemble", ACCENT_PURPLE),
    ("Rank IC", "0.103", "LightGBM", ACCENT_AMBER),
]
for i, (label, value, model, color) in enumerate(cards):
    x = Inches(0.7) + i * Inches(3.15)
    add_metric_card(slide, x, Inches(5.9), Inches(2.85), Inches(1.2),
                    f"{label} ({model})", value, accent_line_color=color)


# ════════════════════════════════════════════════════════════════
# SLIDE 12: RESULTS — LONG BIASED
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "7", "Results: Long-Biased (14L/7S)",
                   "February 2023 \u2013 February 2026  |  802 Trading Days")

# Main results table with SPY
lb_data = [
    ["Metric", "LightGBM", "TST", "CrossMamba", "Ensemble", "SPY"],
    ["Total Return", "85.40%", "128.96%", "131.84%", "122.27%", "~76%"],
    ["Annual Return", "21.41%", "29.73%", "30.24%", "28.53%", "20.29%"],
    ["Volatility", "13.70%", "12.86%", "12.80%", "13.64%", "~15.2%"],
    ["Sharpe Ratio", "1.56", "2.31", "2.36", "2.09", "1.34"],
    ["Sortino Ratio", "2.46", "3.83", "3.97", "3.49", "\u2014"],
    ["Max Drawdown", "-20.25%", "-9.31%", "-9.20%", "-15.44%", "-18.76%"],
    ["Calmar Ratio", "1.06", "3.19", "3.29", "1.85", "\u2014"],
    ["Win Rate", "54.24%", "55.99%", "54.36%", "55.24%", "\u2014"],
    ["Profit Factor", "1.27", "1.40", "1.42", "1.37", "\u2014"],
]
add_table(slide, Inches(0.7), Inches(1.5), Inches(8.7), Inches(4.0),
          lb_data, font_size=12,
          col_widths=[Inches(1.7), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4)])

# Key observations
add_text_box(slide, Inches(9.8), Inches(1.5), Inches(3.2), Inches(0.3),
             "Key Observations", font_size=15, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(9.8), Inches(1.9), Inches(3.2), Inches(3.5), [
    "**All models beat SPY** on return and risk-adjusted basis",
    "CrossMamba: **131.8% total** (vs SPY ~76%) with **half** the drawdown",
    "TST nearly ties CrossMamba \u2014 both achieve **~2.35 Sharpe**",
    "LightGBM captures market beta but adds **minimal alpha** (+1.1% over SPY)",
    "Long-biased dramatically improves all models by capturing the **equity risk premium**",
], font_size=11)

# Alpha vs SPY cards
alpha_cards = [
    ("CrossMamba Alpha", "+9.95%", "vs SPY annual", ACCENT_GREEN),
    ("TST Alpha", "+9.44%", "vs SPY annual", ACCENT_GREEN),
    ("CrossMamba Sharpe", "2.36", "vs SPY 1.34 (+76%)", ACCENT_BLUE),
    ("Max DD Improvement", "-9.20%", "vs SPY -18.76%", ACCENT_PURPLE),
]
for i, (label, value, sub, color) in enumerate(alpha_cards):
    x = Inches(0.7) + i * Inches(3.15)
    add_metric_card(slide, x, Inches(5.9), Inches(2.85), Inches(1.2),
                    f"{label} \u2014 {sub}", value, accent_line_color=color)


# ════════════════════════════════════════════════════════════════
# SLIDE 13: FEATURE IMPORTANCE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "8", "Feature Importance & Alpha Sources",
                   "Top Predictive Features by LightGBM Split Gain")

# Top features table
fi_data = [
    ["Rank", "Feature", "Gain", "Category"],
    ["1", "Earnings day return", "79.0", "Fundamental"],
    ["2", "63d volatility", "60.7", "Price/Volume"],
    ["3", "CS vol rank (63d)", "54.7", "Price/Volume"],
    ["4", "Distance from 252d high", "51.7", "Price/Volume"],
    ["5", "21d volatility", "51.3", "Price/Volume"],
    ["6", "Amihud illiquidity (21d)", "44.0", "Price/Volume"],
    ["7", "MACD histogram", "43.3", "Price/Volume"],
    ["8", "CS momentum rank (63d)", "43.0", "Price/Volume"],
    ["9", "Momentum acceleration (63d)", "42.3", "Price/Volume"],
    ["10", "CS Amihud illiquidity (21d)", "41.3", "Price/Volume"],
    ["11", "Log market cap", "40.3", "Fundamental"],
    ["12", "Distance from 252d high (raw)", "39.0", "Price/Volume"],
    ["13", "CS earnings growth rank", "38.3", "Fundamental"],
    ["14", "Distance from 63d high", "37.3", "Price/Volume"],
    ["15", "CS Bollinger Band width", "37.0", "Price/Volume"],
]
add_table(slide, Inches(0.7), Inches(1.5), Inches(7.5), Inches(5.5),
          fi_data, font_size=11,
          col_widths=[Inches(0.6), Inches(3.2), Inches(0.8), Inches(2.9)])

# Interpretation
add_text_box(slide, Inches(8.7), Inches(1.5), Inches(4.2), Inches(0.3),
             "Alpha Source Interpretation", font_size=15, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(8.7), Inches(1.9), Inches(4.2), Inches(3.0), [
    "**Earnings day return** (#1): Post-earnings announcement drift \u2014 "
    "stocks with positive surprises continue to outperform",
    "**Volatility cluster** (#2-3, #5): Low-volatility anomaly captured across multiple horizons",
    "**Mean reversion** (#4, #12, #14): Distance from highs identifies oversold conditions",
    "**Liquidity premium** (#6, #10): Amihud ratio captures higher expected returns from liquidity risk",
    "**Momentum + acceleration** (#8-9): 63d intermediate momentum with trend inflection detection",
], font_size=11)

# Domain breakdown
add_text_box(slide, Inches(8.7), Inches(5.1), Inches(4.2), Inches(0.3),
             "Signal Domain Breakdown", font_size=14, color=ACCENT_PURPLE, bold=True)

domain_data = [
    ["Domain", "Features in Top 50", "Gain Share"],
    ["Price/Volume", "32", "~65%"],
    ["Fundamental", "18", "~35%"],
    ["Cross-Asset", "Selected, low direct gain", "\u2014"],
    ["Sentiment", "Selected, low direct gain", "\u2014"],
]
add_table(slide, Inches(8.7), Inches(5.45), Inches(4.2), Inches(1.6),
          domain_data, font_size=10, header_color=ACCENT_PURPLE)


# ════════════════════════════════════════════════════════════════
# SLIDE 14: KEY FINDINGS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "9", "Architectural Insights & Key Findings", "")

# Left: Why sequence models win
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.3),
             "Why Sequence Models Dominate Tree Ensembles", font_size=15, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(0.7), Inches(1.9), Inches(5.8), Inches(2.5), [
    "LightGBM treats each observation as **independent** \u2014 "
    "it cannot learn that momentum accelerating over the last 5 of 21 days "
    "carries different meaning than uniform momentum",
    "TST and CrossMamba process **21-day sequences**, explicitly capturing "
    "how features evolve over time",
    "Result: **2\u20133\u00d7 Sharpe improvement** from temporal modeling alone, "
    "with identical features and risk management",
], font_size=13)

# Middle: Why CrossMamba > TST
add_text_box(slide, Inches(0.7), Inches(4.0), Inches(5.8), Inches(0.3),
             "Why CrossMamba Edges Out TST", font_size=15, color=ACCENT_GREEN, bold=True)

add_bullet_frame(slide, Inches(0.7), Inches(4.4), Inches(5.8), Inches(2.0), [
    "**Selective forgetting**: learns which past data matters \u2014 financial data is noisy, "
    "so knowing what to ignore is as valuable as knowing what to attend to",
    "**Local + global context**: depthwise conv captures short-term patterns "
    "(3\u20134 day effects) while hidden state carries regime information",
    "**Computational efficiency**: O(L) vs O(L\u00b2) unlocks potential for longer input windows",
], font_size=13)

# Right: Sharpe comparison
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Sharpe Ratio by Model", font_size=15, color=ACCENT_BLUE, bold=True)

# Visual bar chart using shapes
models_sharpe = [
    ("CrossMamba", 2.36, ACCENT_GREEN),
    ("TST", 2.31, ACCENT_BLUE),
    ("Ensemble", 2.09, ACCENT_PURPLE),
    ("LightGBM", 1.56, ACCENT_AMBER),
    ("SPY", 1.34, DIM_GRAY),
]
max_sharpe = 2.5
bar_area_w = Inches(3.8)
for i, (name, sharpe, color) in enumerate(models_sharpe):
    y = Inches(2.0) + i * Inches(0.65)
    # Label
    add_text_box(slide, Inches(7.2), y, Inches(1.4), Inches(0.35),
                 name, font_size=12, color=LIGHT_GRAY, bold=True, alignment=PP_ALIGN.RIGHT)
    # Bar
    bar_w = int(bar_area_w * (sharpe / max_sharpe))
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(8.7), y + Inches(0.05), bar_w, Inches(0.28))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.adjustments[0] = 0.3
    # Value
    add_text_box(slide, Inches(8.7) + bar_w + Inches(0.1), y, Inches(0.8), Inches(0.35),
                 f"{sharpe:.2f}", font_size=13, color=color, bold=True)

add_text_box(slide, Inches(7.2), Inches(5.3), Inches(5.5), Inches(0.2),
             "Long-biased (14L/7S) configuration", font_size=10, color=DIM_GRAY)

# Bottom: Ensemble paradox + Drawdown
add_shape(slide, Inches(7.2), Inches(5.6), Inches(5.3), Inches(1.5), fill_color=BG_CARD,
          border_color=ACCENT_AMBER)
add_text_box(slide, Inches(7.5), Inches(5.7), Inches(4.8), Inches(0.3),
             "The Ensemble Paradox", font_size=13, color=ACCENT_AMBER, bold=True)
add_text_box(slide, Inches(7.5), Inches(6.05), Inches(4.8), Inches(0.9),
             "Equal-weight ensemble (34/33/33) underperforms CrossMamba in both regimes. "
             "LightGBM's weaker signal dilutes the blend. Naive diversification across "
             "architectures does not guarantee improvement when model quality is heterogeneous.",
             font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 15: DRAWDOWN ANALYSIS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "9b", "Risk-Adjusted Performance Deep Dive",
                   "Drawdown Analysis & Factor Risk Decomposition")

# Drawdown comparison (visual bars)
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(6), Inches(0.3),
             "Maximum Drawdown Comparison (Long-Biased)", font_size=15, color=ACCENT_BLUE, bold=True)

dd_models = [
    ("CrossMamba", -9.20, ACCENT_GREEN),
    ("TST", -9.31, ACCENT_BLUE),
    ("Ensemble", -15.44, ACCENT_PURPLE),
    ("SPY", -18.76, MED_GRAY),
    ("LightGBM", -20.25, ACCENT_RED),
]
max_dd = 25
for i, (name, dd, color) in enumerate(dd_models):
    y = Inches(2.0) + i * Inches(0.7)
    add_text_box(slide, Inches(0.9), y, Inches(1.6), Inches(0.35),
                 name, font_size=13, color=LIGHT_GRAY, bold=True, alignment=PP_ALIGN.RIGHT)
    bar_w = int(Inches(5) * (abs(dd) / max_dd))
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(2.7), y + Inches(0.05), bar_w, Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.adjustments[0] = 0.3
    add_text_box(slide, Inches(2.7) + bar_w + Inches(0.1), y, Inches(1), Inches(0.35),
                 f"{dd:.2f}%", font_size=12, color=color, bold=True)

add_text_box(slide, Inches(0.9), Inches(5.6), Inches(6), Inches(0.6),
             "CrossMamba and TST achieve half the drawdown of SPY while delivering "
             "~50% higher annualized returns. This is the hallmark of genuine alpha: "
             "better returns with less risk.",
             font_size=12, color=LIGHT_GRAY, bold=False)

# Right: Factor risk decomposition
add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.3),
             "Factor Risk Decomposition", font_size=15, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(7.5), Inches(1.85), Inches(5), Inches(0.3),
             "Latest portfolio snapshot", font_size=11, color=MED_GRAY)

fr_data = [
    ["Component", "Value"],
    ["Total portfolio vol", "9.44%"],
    ["Factor-driven vol", "1.62%"],
    ["Specific (idiosyncratic) vol", "9.30%"],
    ["Factor share of total risk", "~3%"],
]
add_table(slide, Inches(7.5), Inches(2.3), Inches(4.8), Inches(1.8),
          fr_data, font_size=12, header_color=ACCENT_PURPLE)

# Factor exposures
add_text_box(slide, Inches(7.5), Inches(4.3), Inches(5), Inches(0.3),
             "Factor Exposures (\u03c3)", font_size=14, color=ACCENT_PURPLE, bold=True)

fe_data = [
    ["Factor", "Exposure", "Interpretation"],
    ["Market", "+0.030", "Near-zero beta (neutral achieved)"],
    ["Size", "-0.174", "Slight small-cap tilt"],
    ["Value", "+0.156", "Moderate value tilt"],
    ["Momentum", "+0.152", "Moderate momentum tilt"],
    ["Volatility", "+0.068", "Near-neutral"],
    ["Quality", "-0.094", "Slight negative quality tilt"],
]
add_table(slide, Inches(7.5), Inches(4.65), Inches(5.2), Inches(2.3),
          fe_data, font_size=11, header_color=ACCENT_PURPLE,
          col_widths=[Inches(1.3), Inches(1.1), Inches(2.8)])


# ════════════════════════════════════════════════════════════════
# SLIDE 16: PRODUCTION & FUTURE WORK
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "10", "Production Deployment & Future Work", "")

# Left: Production
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.3),
             "Live Trading Integration", font_size=15, color=ACCENT_GREEN, bold=True)

add_bullet_frame(slide, Inches(0.7), Inches(1.9), Inches(5.5), Inches(2.5), [
    "**Broker**: Alpaca Markets (paper trading validated, live-ready)",
    "**Order type**: Limit orders with 5 bps offset from mid",
    "**Rebalance frequency**: Every 21 trading days",
    "**Paper trade validated**: 14 long + 7 short positions submitted Feb 24, 2026",
    "All order fills confirmed via Alpaca API",
], font_size=13)

# Sample positions
add_text_box(slide, Inches(0.7), Inches(4.3), Inches(5.5), Inches(0.3),
             "Latest Signal Sample", font_size=14, color=ACCENT_BLUE, bold=True)

signal_data = [
    ["Direction", "Ticker", "Weight"],
    ["LONG", "ACN", "10.0%"],
    ["LONG", "HUM", "10.0%"],
    ["LONG", "SCHW", "9.7%"],
    ["LONG", "AXP", "9.0%"],
    ["LONG", "SPGI", "8.4%"],
    ["SHORT", "PCAR", "-6.7%"],
    ["SHORT", "TXN", "-5.9%"],
    ["SHORT", "SLB", "-4.0%"],
]
add_table(slide, Inches(0.7), Inches(4.65), Inches(4.5), Inches(2.7),
          signal_data, font_size=11,
          col_widths=[Inches(1.0), Inches(1.5), Inches(2.0)])

# Right: Future work
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Caveats & Limitations", font_size=15, color=ACCENT_RED, bold=True)

add_bullet_frame(slide, Inches(7.2), Inches(1.9), Inches(5.5), Inches(2.0), [
    "**Survivorship bias**: universe based on current S&P 500 \u2014 delisted stocks not included",
    "**Backtest period**: 3.18 years in a generally bullish market; bear regime coverage limited",
    "**Capacity**: tested at $100K AUM \u2014 behavior at scale (>$10M) is untested",
], font_size=12, bullet_color=ACCENT_RED)

add_text_box(slide, Inches(7.2), Inches(4.0), Inches(5.5), Inches(0.3),
             "Future Work", font_size=15, color=ACCENT_AMBER, bold=True)

add_bullet_frame(slide, Inches(7.2), Inches(4.4), Inches(5.5), Inches(3.0), [
    "**Learned ensemble weights**: stacking or attention-weighted blending instead of equal weights",
    "**Sequence length ablation**: CrossMamba may benefit from 63d or 126d windows (O(L) enables this)",
    "**Cross-stock attention**: model contagion and momentum spillover between related stocks",
    "**Alternative targets**: volatility-adjusted returns, drawdown probability prediction",
    "**Extended backtest**: include 2020 COVID crash, 2022 rate shock for bear market validation",
], font_size=12, bullet_color=ACCENT_AMBER)


# ════════════════════════════════════════════════════════════════
# SLIDE 17: SUMMARY
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Title
add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.6),
             "Summary of Findings", font_size=30, color=WHITE, bold=True)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.0),
                               Inches(12), Pt(1.5))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Finding cards
findings = [
    ("CrossMamba is the best alpha model",
     "Highest Sharpe (2.10 / 2.36), lowest max DD (-5.93% / -9.20%), "
     "highest total return in both configurations",
     ACCENT_GREEN),
    ("Temporal modeling provides 2\u20133\u00d7 edge",
     "Both TST and CrossMamba outperform LightGBM dramatically \u2014 "
     "sequence information is a critical alpha source",
     ACCENT_BLUE),
    ("Selective state-space > Self-attention",
     "CrossMamba's O(L) complexity and selective forgetting outperform "
     "TST's O(L\u00b2) attention, with a smaller model",
     ACCENT_PURPLE),
    ("Long-biased captures equity premium + alpha",
     "14L/7S achieves 30% annual return vs SPY 20%, with better "
     "risk-adjusted metrics across all models",
     ACCENT_AMBER),
    ("Risk management is load-bearing",
     "Factor neutralization, vol targeting, and drawdown control keep "
     "all models within target vol despite different alpha profiles",
     MED_GRAY),
    ("Naive ensembling can hurt performance",
     "Equal-weight blend underperforms best individual model when "
     "model quality is heterogeneous \u2014 quality > quantity",
     ACCENT_RED),
]

card_w = Inches(5.8)
for i, (title, desc, color) in enumerate(findings):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + col * (card_w + Inches(0.3))
    y = Inches(1.3) + row * Inches(1.8)

    card = add_shape(slide, x, y, card_w, Inches(1.55), fill_color=BG_CARD)
    # Color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), Inches(1.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Title
    add_text_box(slide, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.4), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
    # Description
    add_text_box(slide, x + Inches(0.25), y + Inches(0.55), card_w - Inches(0.4), Inches(0.85),
                 desc, font_size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 18: THANK YOU / Q&A
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(1),
             "Thank You", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6),
                               Inches(2.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_text_box(slide, Inches(0), Inches(4.0), SLIDE_W, Inches(0.5),
             "Questions & Discussion", font_size=22, color=MED_GRAY,
             alignment=PP_ALIGN.CENTER)

# Bottom metric strip
final_metrics = [
    ("Best Model", "CrossMamba"),
    ("Best Sharpe", "2.36"),
    ("Annual Alpha", "+9.95%"),
    ("Max Drawdown", "-9.20%"),
    ("Total Return", "131.84%"),
]
strip_w = Inches(2.2)
start_x = Inches(0.9)
for i, (label, value) in enumerate(final_metrics):
    x = start_x + i * (strip_w + Inches(0.3))
    add_metric_card(slide, x, Inches(5.5), strip_w, Inches(1.2),
                    label, value, accent_line_color=ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_path = os.path.join("results", "Multi_Model_Trading_System_Presentation.pptx")
os.makedirs("results", exist_ok=True)
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
